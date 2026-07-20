from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any

from sqlalchemy.orm import Session

from browser.net_guard import assert_public_url
from browser.pool import browser_pool
from config.config import workspace_root
from db.models import WorkspaceFile
from rag.retrieve import search_knowledge
from tools.web_tools import validate_and_normalize_args, web_fetch, web_search
from utils.logger import get_logger

logger = get_logger("tools")

CONFIRM_TOOLS = {
    "browser_submit",
    "http_request_write",
    "file_delete",
}


def user_workspace(user_id: int) -> Path:
    path = workspace_root / str(user_id)
    path.mkdir(parents=True, exist_ok=True)
    return path


def register_workspace_file(
    db: Session,
    *,
    user_id: int,
    conversation_id: int | None,
    path: Path,
    mime: str,
) -> WorkspaceFile:
    row = WorkspaceFile(
        user_id=user_id,
        conversation_id=conversation_id,
        path=str(path),
        name=path.name,
        mime=mime,
        size=path.stat().st_size if path.exists() else 0,
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return row


async def execute_tool(
    name: str,
    args: dict[str, Any],
    *,
    db: Session,
    user_id: int,
    conversation_id: int,
) -> dict[str, Any]:
    # OpenClaw 风格：执行前校验/归一化入参
    normalized, err = validate_and_normalize_args(name, args)
    if err:
        return {"error": err}
    if normalized is not None:
        args = normalized

    if name == "web_search":
        return await web_search(str(args.get("query") or ""), user_id=user_id)
    if name == "web_fetch":
        return await web_fetch(str(args.get("url") or ""))

    if name == "browser_extract":
        selector = args.get("selector")
        # 兼容模型乱写的 targets/css 格式
        if not selector and isinstance(args.get("targets"), list) and args["targets"]:
            t0 = args["targets"][0]
            if isinstance(t0, dict):
                selector = t0.get("css") or t0.get("selector")
            elif isinstance(t0, str):
                selector = t0
        if not selector and args.get("css"):
            selector = args.get("css")
        # 先看当前页；about:blank 直接失败，由编排层决定是否自动 navigate
        peek = await browser_pool.peek(user_id)
        url = str(peek.get("url") or "about:blank")
        if url in {"", "about:blank"}:
            return {
                "error": "浏览器尚未打开页面(about:blank)，请先 browser_navigate 或改用 web_fetch/web_search",
                "url": url,
            }
        result = await browser_pool.extract(user_id, selector if isinstance(selector, str) else None)
        return result
    if name == "browser_screenshot":
        peek = await browser_pool.peek(user_id)
        url = str(peek.get("url") or "about:blank")
        if url in {"", "about:blank"}:
            return {"error": "浏览器尚未打开页面，请先 browser_navigate", "url": url}
        return await browser_pool.screenshot(user_id)
    if name == "browser_click":
        peek = await browser_pool.peek(user_id)
        if str(peek.get("url") or "about:blank") in {"", "about:blank"}:
            return {"error": "浏览器尚未打开页面，请先 browser_navigate"}
        return await browser_pool.click(user_id, str(args.get("selector") or ""))
    if name == "browser_type":
        peek = await browser_pool.peek(user_id)
        if str(peek.get("url") or "about:blank") in {"", "about:blank"}:
            return {"error": "浏览器尚未打开页面，请先 browser_navigate"}
        return await browser_pool.type_text(
            user_id, str(args.get("selector") or ""), str(args.get("text") or "")
        )
    if name == "browser_submit":
        # 确认后执行：点击提交按钮
        peek = await browser_pool.peek(user_id)
        if str(peek.get("url") or "about:blank") in {"", "about:blank"}:
            return {"error": "浏览器尚未打开页面，请先 browser_navigate"}
        return await browser_pool.click(user_id, str(args.get("selector") or "button[type=submit]"))
    if name == "browser_navigate":
        return await browser_pool.navigate(user_id, str(args.get("url") or ""))
    if name == "kb_search":
        hits = await search_knowledge(db, user_id, str(args.get("query") or ""), top_k=5)
        return {"hits": hits}
    if name == "http_request":
        method = str(args.get("method") or "GET").upper()
        try:
            url = assert_public_url(str(args.get("url") or ""))
        except Exception as exc:
            return {"error": str(exc)}
        if method not in {"GET", "HEAD", "OPTIONS"}:
            return {"error": "写操作请使用 http_request_write 并经确认"}
        # 高层推荐：对 HTML 页优先提示 web_fetch（仍执行以兼容）
        fetched = await web_fetch(url)
        if fetched.get("error"):
            return fetched
        return {
            "status": fetched.get("status"),
            "url": fetched.get("url"),
            "text": fetched.get("text"),
        }
    if name == "http_request_write":
        method = str(args.get("method") or "POST").upper()
        url = assert_public_url(str(args.get("url") or ""))
        import httpx

        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            resp = await client.request(
                method,
                url,
                json=args.get("json"),
                content=args.get("body"),
                headers=args.get("headers") or {},
            )
        return {"status": resp.status_code, "text": resp.text[:8000]}
    if name == "file_list":
        root = user_workspace(user_id)
        files = [p.name for p in root.iterdir() if p.is_file()]
        return {"files": files}
    if name == "file_write_markdown":
        return _write_text_file(db, user_id, conversation_id, args, ".md", "text/markdown")
    if name == "file_write_html":
        return _write_text_file(db, user_id, conversation_id, args, ".html", "text/html")
    if name == "file_write_docx":
        return _write_docx(db, user_id, conversation_id, args)
    if name == "file_write_xlsx":
        return _write_xlsx(db, user_id, conversation_id, args)
    if name == "file_write_pptx":
        return _write_pptx(db, user_id, conversation_id, args)
    if name == "file_write_pdf":
        return _write_pdf(db, user_id, conversation_id, args)
    if name == "file_delete":
        name_ = str(args.get("name") or "")
        path = user_workspace(user_id) / Path(name_).name
        if path.exists():
            path.unlink()
        return {"deleted": name_}
    if name == "run_code":
        from tools.code_sandbox import run_python_sandbox

        code = str(args.get("code") or args.get("source") or args.get("python") or "")
        timeout = int(args.get("timeout") or 8)
        return run_python_sandbox(code, workdir=user_workspace(user_id), timeout_sec=timeout)
    if name == "memory_recall":
        from agent.session_memory import recall_session_memory

        return recall_session_memory(
            db,
            user_id=user_id,
            conversation_id=conversation_id,
            summary_id=str(args.get("summary_id") or ""),
            query=str(args.get("query") or ""),
        )
    if name in {"feishu_search", "meeting_search"}:
        return {"error": "该连接器未启用（一期预留）"}
    return {"error": f"未知工具: {name}"}


def _write_text_file(
    db: Session,
    user_id: int,
    conversation_id: int,
    args: dict,
    suffix: str,
    mime: str,
) -> dict:
    from tools.web_tools import _pick_write_content

    filename = str(args.get("filename") or f"note{suffix}")
    if not filename.endswith(suffix):
        filename += suffix
    content = _pick_write_content(args) or str(args.get("content") or "")
    if len(content.strip()) < 8:
        return {"ok": False, "error": "正文为空，未写入文件。请在 content 中提供完整内容后重试。"}
    path = user_workspace(user_id) / Path(filename).name
    path.write_text(content, encoding="utf-8")
    row = register_workspace_file(
        db, user_id=user_id, conversation_id=conversation_id, path=path, mime=mime
    )
    return {"file_id": row.id, "name": row.name, "path": row.path, "chars": len(content)}


def _write_docx(db: Session, user_id: int, conversation_id: int, args: dict) -> dict:
    from docx import Document
    from tools.web_tools import _pick_write_content

    filename = str(args.get("filename") or "document.docx")
    if not filename.endswith(".docx"):
        filename += ".docx"
    content = _pick_write_content(args) or str(args.get("content") or "")
    if len(content.strip()) < 8:
        return {"ok": False, "error": "正文为空，未生成 Word。请把完整内容写入 content 后重试。"}
    path = user_workspace(user_id) / Path(filename).name
    doc = Document()
    title = args.get("title")
    if title:
        doc.add_heading(str(title), 0)
    for para in content.split("\n"):
        line = para.rstrip()
        if not line.strip():
            doc.add_paragraph("")
            continue
        stripped = line.strip()
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=1)
        else:
            doc.add_paragraph(stripped)
    doc.save(path)
    # 二次校验：文件过小视为失败
    size = path.stat().st_size if path.exists() else 0
    if size < 1200 and len(content) < 20:
        return {"ok": False, "error": "生成的 Word 几乎为空，请重试并提供更长正文。"}
    row = register_workspace_file(
        db,
        user_id=user_id,
        conversation_id=conversation_id,
        path=path,
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    return {"file_id": row.id, "name": row.name, "chars": len(content), "size": size}


def _write_xlsx(db: Session, user_id: int, conversation_id: int, args: dict) -> dict:
    from openpyxl import Workbook

    filename = str(args.get("filename") or "sheet.xlsx")
    if not filename.endswith(".xlsx"):
        filename += ".xlsx"
    path = user_workspace(user_id) / Path(filename).name
    wb = Workbook()
    ws = wb.active
    ws.title = str(args.get("sheet_name") or "Sheet1")
    rows = args.get("rows") or []
    if isinstance(rows, str):
        try:
            rows = json.loads(rows)
        except Exception:
            rows = [[rows]]
    for r in rows:
        if isinstance(r, (list, tuple)):
            ws.append(list(r))
        else:
            ws.append([r])
    wb.save(path)
    row = register_workspace_file(
        db,
        user_id=user_id,
        conversation_id=conversation_id,
        path=path,
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    return {"file_id": row.id, "name": row.name}


def _write_pptx(db: Session, user_id: int, conversation_id: int, args: dict) -> dict:
    from pptx import Presentation

    filename = str(args.get("filename") or "slides.pptx")
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    path = user_workspace(user_id) / Path(filename).name
    prs = Presentation()
    slides = args.get("slides") or [{"title": args.get("title") or "演示", "body": args.get("content") or ""}]
    if isinstance(slides, str):
        try:
            slides = json.loads(slides)
        except Exception:
            slides = [{"title": "演示", "body": slides}]
    for item in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = str(item.get("title") or "")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = str(item.get("body") or "")
    prs.save(path)
    row = register_workspace_file(
        db,
        user_id=user_id,
        conversation_id=conversation_id,
        path=path,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return {"file_id": row.id, "name": row.name}


def _write_pdf(db: Session, user_id: int, conversation_id: int, args: dict) -> dict:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from tools.web_tools import _pick_write_content

    filename = str(args.get("filename") or "document.pdf")
    if not filename.endswith(".pdf"):
        filename += ".pdf"
    path = user_workspace(user_id) / Path(filename).name
    text = _pick_write_content(args) or str(args.get("content") or args.get("title") or "")
    if len(text.strip()) < 8:
        return {"ok": False, "error": "正文为空，未生成 PDF。请提供 content。"}
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4
    y = height - 50
    # 使用内置字体，中文可能显示为方框；后续可挂载系统字体
    for line in text.split("\n"):
        c.drawString(40, y, line[:90])
        y -= 16
        if y < 40:
            c.showPage()
            y = height - 50
    c.save()
    row = register_workspace_file(
        db, user_id=user_id, conversation_id=conversation_id, path=path, mime="application/pdf"
    )
    return {
        "file_id": row.id,
        "name": row.name,
        "chars": len(text),
        "note": "PDF 使用内置西文字体，中文建议用 docx",
    }
